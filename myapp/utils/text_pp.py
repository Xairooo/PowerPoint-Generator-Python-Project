import io
import json
import os

import requests
import time
import datetime
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from urllib.parse import quote_plus
from dotenv import load_dotenv
from slide.title_slide import TitleSlide
from slide.grid import GridSlide
from slide.title_content_bullets import TitleContentWithBullets
from slide.arrow_shape import ArrowShape
from slide.comparison import ComparisonSlide
from slide.content_with_caption import ContentWithCaptionSlide
from slide.picture_with_caption import PictureWithCaptionSlide
from slide.table import TableSlide
from slide.TitleAndVerticalTextSlide import TitleAndVerticalTextSlide
from slide.VerticalTitleAndTextSlide import VerticalTitleAndTextSlide
from slide.title_content import TitleContentSlide
from slide.two_columns import TwoColumnSlide
dir_path = 'static/presentations'

load_dotenv()
API_KEY = os.getenv('PEXELS_API_KEY')

def parse_response(response):
    slides = response.split('\n\n')
    slides_content = []
    for slide in slides:
        lines = slide.split('\n')
        title_line = lines[0]
        if ': ' in title_line:
            title = title_line.split(': ', 1)[1]  # Extract the title after 'Slide X: '
        else:
            title = title_line
        content_lines = [line for line in lines[1:] if line != 'Content:']  # Skip line if it is 'Content:'
        content = '\n'.join(content_lines)  # Join the lines to form the content
        # Extract the keyword from the line that starts with 'Keyword:'
        keyword_line = [line for line in lines if 'Keyword:' or 'Keywords:' in line][0]
        keyword = keyword_line.split(': ', 1)[1]
        slides_content.append({'title': title, 'content': content, 'keyword': keyword})
    return slides_content


def search_pexels_images(keyword):
    query = quote_plus(keyword.lower())
    print("Query:", query) # Debug
    PEXELS_API_URL = f'https://api.pexels.com/v1/search?query={query}&per_page=1'
    print("URL:", PEXELS_API_URL) # Debug
    headers = {
        'Authorization': API_KEY
    }
    response = requests.get(PEXELS_API_URL, headers=headers)
    print("Response Status Code:", response.status_code) # Debug
    print("Response Content:", response.text) # Debug
    data = json.loads(response.text)
    if 'photos' in data:
        if len(data['photos']) > 0:
            return data['photos'][0]['src']['medium']
    return None


def delete_first_two_slides(presentation):
    slide_ids = [1, 0]
    for slide_id in slide_ids:
        if slide_id < len(presentation.slides):
            xml_slides = presentation.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[slide_id])


def create_ppt(slides_content, template_choice, presentation_title, presenter_name, insert_image):
    # template_path = os.path.join(dir_path, f"{template_choice}.pptx")

    print(template_choice)

    prs = Presentation(template_choice)
    template_choice = 'dark_modern'

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    print(title_slide_layout, content_slide_layout)

    # add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    print(slide.shapes.title)
    title.text = presentation_title

    #add subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = f"Presented by {presenter_name}"

    if template_choice == 'dark_modern':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Times New Roman'
                run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color

    elif template_choice == 'bright_modern':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Arial'
                run.font.color.rgb = RGBColor(255, 20, 147)  # RGB for deep pink color

    # add content slides
    for slide_content in slides_content:
        slide = prs.slides.add_slide(content_slide_layout)

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = slide_content['title']
                if template_choice == 'dark_modern':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = slide_content['content']
                if template_choice == 'dark_modern':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 255, 255)  # RGB for white color

        if insert_image:
            # fetch image URL from Pixabay based on the slide's title
            image_url = search_pexels_images(slide_content['keyword'])
            print("Image URL:", image_url) #debug
            if image_url is not None:
                # download the image
                image_data = requests.get(image_url).content
                # load image into BytesIO object
                image_stream = io.BytesIO(image_data)
                # add the image at the specified position
                slide_width = Inches(20)
                slide_height = Inches(15)

                image_width = Inches(8)  # width of image
                image_height = Inches(5)  # height of image

                left = slide_width - image_width  # calculate left position
                top = slide_height - image_height - Inches(4)  # calculate top position

                slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)

    # add credits slide
    slide = prs.slides.add_slide(content_slide_layout)
    if template_choice == 'dark_modern':
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Times New Roman'
                        run.font.color.rgb = RGBColor(255, 165, 0)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Times New Roman'
                        run.font.color.rgb = RGBColor(255, 255, 255)

    elif template_choice == 'bright_modern':
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(255, 20, 147)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)

    else:
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)

    # Delete the first two slides after all new slides have been added
    delete_first_two_slides(prs)

    # Save the presentation
    prs.save(os.path.join('generated', f"generated_presentation_{datetime.datetime.now()}.pptx"))

def create_ppt_v2(slides_content, template_choice, presentation_title, presenter_name, insert_image):
    # template_path = os.path.join(dir_path, f"{template_choice}.pptx")

    prs = Presentation(template_choice)
    template_choice = 'dark_modern'

    # xml_slides = prs.slides._sldIdLst
    # for _ in range(len(prs.slides)):
    #     xml_slides.remove(xml_slides[0])

    TitleSlide(prs, presentation_title, f"Presented by {presenter_name}")
    content_slide_layout = prs.slide_layouts[1]

    TitleContentWithBullets(prs, "Adding a Bullet Slide", [
        ("Find the bullet slide layout", 0),
        ("Use _TextFrame.text for first bullet", 1),
        ("Use _TextFrame.add_paragraph() for subsequent bullets", 2)
    ])

    GridSlide(
        prs,
        "Grid Slide Example", 
        rows=3, 
        cols=4, 
        cell_width=Inches(1.5), 
        cell_height=Inches(1), 
        spacing=Inches(0.2)
    )

    ArrowShape(prs, "Xin chao vietnam")

    ComparisonSlide(
        prs,
        "Comparison Slide Example",
        "Left Column Title",
        "This is the left column content.",
        "Right Column Title",
        "This is the right column content."
    )

    ContentWithCaptionSlide(
        prs,
        "Left Column Title",
        "This is the left column content.",
        "Right Column Title",
    )

    PictureWithCaptionSlide(
        prs,
        "https://img.freepik.com/premium-photo/wallpaper-with-dark-dramatic-gradient-colors-ai-generated_88211-6704.jpg",
        "This is the picture with caption.",
        "This is the caption.",
    )

    TableSlide(prs, "Adding a Table")

    TitleAndVerticalTextSlide(prs, "Adding a Title and Vertical Text", "This is the vertical text.")

    # TwoColumnSlide(prs, "Adding a Two Column Slide", "This is the left column content.", "This is the right column content.")

    with open('sample_response_1.json') as f:
        data = json.load(f)
        slide_content = data['slides']

    for slide in slide_content:
        title = slide['title']
        layout = slide['layout']
        if layout['type'] == 'TITLE_AND_CONTENT':
            TitleContentSlide.from_json(prs, slide)
        elif layout['type'] == 'TWO_CONTENT':
            TwoColumnSlide.from_json(prs, slide)
            # TwoColumnSlide(prs, title, "'\n'.join(layout['sections'][0]['items'])", "'\n'.join(layout['sections'][1]['content'])")
        elif layout['type'] == 'CONTENT_WITH_CAPTION':
            ContentWithCaptionSlide.from_json(prs, slide)
            # pass
        elif layout['type'] == 'PICTURE_WITH_CAPTION':
            PictureWithCaptionSlide.from_json(prs, slide)
            # pass
        elif layout['type'] == 'COMPARISON':
            ComparisonSlide.from_json(prs, slide)
            # pass
        elif layout['type'] == 'VERTICAL_TITLE_AND_TEXT':
            VerticalTitleAndTextSlide.from_json(prs, slide)
            # pass
        elif layout['type'] == 'TITLE_AND_VERTICAL_TEXT':
            TitleAndVerticalTextSlide.from_json(prs, slide)
        elif layout['type'] == 'TITLE_ONLY':
            TitleSlide.from_json(prs, slide)
            # TitleSlide(prs, title, layout['sections'][0]['content'])
 
    # print(data)

    # for idx, layout in enumerate(prs.slide_layouts):
    #     print(f"Layout {idx}: {layout.name}")
    #     for placeholder in layout.placeholders:
    #         print(f"  Placeholder {placeholder.placeholder_format.idx}: {placeholder.name}")

    # if template_choice == 'dark_modern':
    #     for paragraph in title.text_frame.paragraphs:
    #         for run in paragraph.runs:
    #             run.font.name = 'Times New Roman'
    #             run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color

    # elif template_choice == 'bright_modern':
    #     for paragraph in title.text_frame.paragraphs:
    #         for run in paragraph.runs:
    #             run.font.name = 'Arial'
    #             run.font.color.rgb = RGBColor(255, 20, 147)  # RGB for deep pink color

    # add content slides
    for slide_content in slides_content:
        slide = prs.slides.add_slide(content_slide_layout)

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = slide_content['title']
                if template_choice == 'dark_modern':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = slide_content['content']
                if template_choice == 'dark_modern':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 255, 255)  # RGB for white color

        if insert_image:
            # fetch image URL from Pixabay based on the slide's title
            image_url = search_pexels_images(slide_content['keyword'])
            print("Image URL:", image_url) #debug
            if image_url is not None:
                # download the image
                image_data = requests.get(image_url).content
                # load image into BytesIO object
                image_stream = io.BytesIO(image_data)
                # add the image at the specified position
                slide_width = Inches(20)
                slide_height = Inches(15)

                image_width = Inches(8)  # width of image
                image_height = Inches(5)  # height of image

                left = slide_width - image_width  # calculate left position
                top = slide_height - image_height - Inches(4)  # calculate top position

                slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)

    # add credits slide
    slide = prs.slides.add_slide(content_slide_layout)
    if template_choice == 'dark_modern':
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Times New Roman'
                        run.font.color.rgb = RGBColor(255, 165, 0)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Times New Roman'
                        run.font.color.rgb = RGBColor(255, 255, 255)

    elif template_choice == 'bright_modern':
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(255, 20, 147)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)

    else:
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)

    # Delete the first two slides after all new slides have been added
    delete_first_two_slides(prs)

    # Save the presentation
    prs.save(os.path.join('generated', f"generated_presentation_{str(time.time())}.pptx"))
