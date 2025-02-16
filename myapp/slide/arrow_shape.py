from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from .slide_enum import SlideLayout

class ArrowShape:
    def __init__(self, prs, title_text):
        title_only_slide_layout = prs.slide_layouts[SlideLayout.TITLE_AND_CONTENT]  # Title Only layout
        slide = prs.slides.add_slide(title_only_slide_layout)

        # Set title
        slide.shapes.title.text = title_text

        # Add Step 1 (Pentagon)
        left = Inches(0.93)
        top = Inches(3.0)
        width = Inches(1.75)
        height = Inches(1.0)
        shape = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
        shape.text = 'Step 1'
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add Steps 2 through 6 (Chevrons)
        left += width - Inches(0.4)
        width = Inches(2.0)
        for n in range(2, 4):
            shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
            shape.text = f'Step {n}'
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            left += width - Inches(0.4)

    def save(self, filename):
        self.prs.save(filename)