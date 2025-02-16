from pptx import Presentation
from .slide_enum import SlideLayout
class VerticalTitleAndTextSlide:
    def __init__(self, prs, title_text, content_text):
        layout = prs.slide_layouts[SlideLayout.VERTICAL_TITLE_AND_TEXT]  # Vertical Title and Text layout
        self.slide = prs.slides.add_slide(layout)
        self.slide.shapes.title.text = title_text
        self.slide.placeholders[1].text = content_text

    @staticmethod
    def from_json(prs, json_slide):
        title = json_slide["title"]
        sections = json_slide["layout"]["sections"]
        content_text = sections[1]["content"]
        return VerticalTitleAndTextSlide(prs,title,  content_text)
    
    def save(self, filename):
        self.prs.save(filename)