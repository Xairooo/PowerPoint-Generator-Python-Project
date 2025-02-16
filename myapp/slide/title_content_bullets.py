from pptx import Presentation
from .slide_enum import SlideLayout

class TitleContentWithBullets:
    def __init__(self, prs, title_text, bullets):
        bullet_slide_layout = prs.slide_layouts[SlideLayout.TITLE_AND_CONTENT]  # Title and Content layout
        slide = prs.slides.add_slide(bullet_slide_layout)

        # Set title text
        slide.shapes.title.text = title_text

        # for shape in (slide.shapes):
        #     # if shape.is_placeholder:
        #     phf = shape.placeholder_format
        #     print('%d, %s' % (phf.idx, phf.type))

        # Set bullet points
        text_frame = slide.placeholders[1].text_frame
        text_frame.text = bullets[0][0]  # First bullet
        
        for bullet, level in bullets[1:]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = level

    @staticmethod
    def from_json(prs, json_obj):
        title_text = json_obj["title"]
        sections = json_obj["layout"]["sections"]
        main_content = sections[1]
        return TitleContentWithBullets(prs, title_text, bullets)

    def save(self, filename):
        self.prs.save(filename)

# Example Usage
# bullets = [
#     ("Find the bullet slide layout", 0),
#     ("Use _TextFrame.text for first bullet", 1),
#     ("Use _TextFrame.add_paragraph() for subsequent bullets", 2)
# ]
# bullet_slide = TitleContentWithBullets("Adding a Bullet Slide", bullets)
# bullet_slide.save("bullet_slide.pptx")
