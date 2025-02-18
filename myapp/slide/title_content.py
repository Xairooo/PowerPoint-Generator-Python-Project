from pptx import Presentation
from pptx.util import Pt
from .slide_enum import SlideLayout
from utils.common import fallback_string
from models import BulletPointsSection  # Import from your dataclasses  

class TitleContentSlide:
    def __init__(self, prs, title_text, content_section):
        self.prs = prs
        slide_layout = prs.slide_layouts[SlideLayout.TITLE_AND_CONTENT]
        self.slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        self.slide.shapes.title.text = title_text
        
        # Handle different content types
        if isinstance(content_section, BulletPointsSection):
            self._add_bullet_points(content_section.items)
            return

        self.slide.placeholders[1].text = content_section

    def _add_bullet_points(self, items):
        text_frame = self.slide.placeholders[1].text_frame
        text_frame.text = items[0]
        
        for idx, item in enumerate(items[1:]):
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0  # You can modify this based on indentation in your data
            # p.font.size = Pt(18)

    @staticmethod
    def from_json(prs, json_slide):
        title = json_slide["title"]
        layout = json_slide["layout"]
        
        # Find bullet points section
        content_section = [x for x in layout['sections'] if x["type"] == 'text' or x["type"] == 'title'][0]["content"]
        bullets = [x for x in layout['sections'] if x["type"] == 'bullets'][0]['content']
        if content_section and content_section != "":
            return TitleContentSlide(prs, title, fallback_string(content_section))
            
        return TitleContentSlide(prs, title, BulletPointsSection("body", bullets))

    def save(self, filename):
        self.prs.save(filename)