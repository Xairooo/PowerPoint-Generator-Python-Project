from pptx import Presentation
from pptx.util import Inches
from .slide_enum import SlideLayout
from models import BulletPointsSection, ImageSection, TextSection
from utils.common import search_unsplash_images, create_image_stream
from urllib.request import urlretrieve

class TwoColumnSlide:
    def __init__(self, prs, title_text, sections):
        self.prs = prs
        two_content_layout = prs.slide_layouts[SlideLayout.TWO_CONTENT]
        self.slide = prs.slides.add_slide(two_content_layout)
        
        # Set title
        self.slide.shapes.title.text = title_text
        
        # Process sections
        for section in sections:
            placeholder_idx = 0
            if section.position == "left":
                placeholder_idx = 1
            elif section.position == "right":
                placeholder_idx = 2
            else:
                continue  # Skip invalid positions

            # print(placeholder_idx)
            
            if isinstance(section, BulletPointsSection):
                self._add_bullet_points(placeholder_idx, section.items)
            elif isinstance(section, ImageSection):
                self._add_image(placeholder_idx, search_unsplash_images(title_text))
            elif isinstance(section, TextSection):
                self._add_text(placeholder_idx, section.content)

    def _add_bullet_points(self, placeholder_idx, items):
        text_frame = self.slide.placeholders[placeholder_idx].text_frame
        # text_frame.clear()
        text_frame.text = items[0]
        
        for item in items[1:]:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            # p.font.size = Pt(18)

    def _add_image(self, placeholder_idx, image_path):
        placeholder = self.slide.placeholders[placeholder_idx]
        left = placeholder.left
        top = placeholder.top
        # print(left, top)
        self.slide.shapes.add_picture(create_image_stream(image_path), left, top, width=Inches(6), height=Inches(4))
        # Preserve aspect ratio
        # placeholder.width = int(placeholder.width * 0.9)
        # placeholder.height = int(placeholder.height * 0.9)

    def _add_text(self, placeholder_idx, text):
        placeholder = self.slide.placeholders[placeholder_idx]
        if isinstance(text, list):
            placeholder.text = '\n'.join(text)
        else:
            placeholder.text = text

    @staticmethod
    def from_json( prs, json_slide):
        title = json_slide["title"]
        layout = json_slide["layout"]
        
        # Find two-columns section
        # two_col_section = next(
        #     (s for s in layout["sections"] if s["type"] == "two-columns"),
        #     None
        # )
        
        # if not two_col_section:
        #     raise ValueError("No two-columns section found in slide layout")
            
        # Process columns
        sections = []
        for idx, col in enumerate(layout["sections"]):
            position = "left" if idx == 0 else "right"
            if col["type"] == "bullets":
                sections.append(BulletPointsSection(
                    position=position,
                    items=col["content"]
                ))
            elif col["type"] == "image":
                sections.append(ImageSection(
                    position=position,
                    description=col["content"]
                ))
            elif col["type"] == "text":
                sections.append(TextSection(
                    position=position,
                    content=col["content"]
                ))
        return TwoColumnSlide(prs, title, sections)

        

# Example Usage
# two_column_slide = TwoColumnSlide(
#     "Two Column Slide Example",
#     "This is the left column content.",
#     "This is the right column content."
# )
# two_column_slide.save("two_column_slide.pptx")
