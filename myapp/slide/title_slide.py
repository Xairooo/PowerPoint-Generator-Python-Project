from .slide_enum import SlideLayout
from pptx.util import Inches

class TitleSlide:
    def __init__(self, prs, title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[SlideLayout.TITLE])
        # print(self.slide.shapes)
        # print(title_text, subtitle_text)
        if slide.shapes.title:
            # print(" case has title")
            slide.shapes.title.text = title_text
        else:
            # print("case no title")
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            title_box.text_frame.text = title_text

        # print(len(slide.shapes))
        # for shape in (slide.shapes):
        #     # if shape.is_placeholder:
        #     phf = shape.placeholder_format
        #     print('%d, %s' % (phf.idx, phf.type))

        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text 

    def save(self, filename):
        self.prs.save(filename)

    @staticmethod
    def from_json(prs, json_slide):
        title = json_slide["title"]
        sections = json_slide["layout"]["sections"]
        return TitleSlide(prs, title, sections[0]["content"])

    def save(self, filename):
        self.prs.save(filename)
