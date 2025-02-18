from pptx import Presentation
from pptx.util import Inches, Pt
from .slide_enum import SlideLayout
from pptx.enum.shapes import MSO_SHAPE
from utils.common import fallback_string

class ComparisonSlide:
    def __init__(self, prs, title_text, left_title, left_content, right_title, right_content):
        # Find Comparison layout or fallback to Two Content layout
        layout = prs.slide_layouts[SlideLayout.COMPARISON]

        # print(layout.name)

        # Add slide
        slide = prs.slides.add_slide(layout)

        # Set slide title
        if slide.shapes.title:
            slide.shapes.title.text = title_text
        else:
            slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1)).text = title_text

        # Left content placeholder (idx 1) or fallback
        try:
            left_placeholder = slide.placeholders[1]
            left_placeholder.text = left_title
            left_placeholder.text_frame.add_paragraph().text = left_content
        except KeyError:
            slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(3)).text = f"{left_title}\n{left_content}"

        # Right content placeholder (idx 2) or fallback
        try:
            right_placeholder = slide.placeholders[3]
            right_placeholder.text = right_title
            right_placeholder.text_frame.add_paragraph().text = right_content
        except KeyError:
            slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(3)).text = f"{right_title}\n{right_content}"

        # self.add_vertical_separator(slide, Inches(5.0), Inches(1.5), Inches(0.05), Inches(3.5))


    def add_vertical_separator(self, slide, left, top, width, height):
        """Add a vertical line between two columns."""
        line = slide.shapes.add_shape(
            MSO_SHAPE.LINE, left, top, width, height
        )
        line.line.width = Pt(2)
        line.line.color.rgb = (0, 0, 0)  # Black line

    @staticmethod
    def from_json(prs, json_slide):
        title = json_slide["title"]
        sections = json_slide["layout"]["sections"]
        left_title = sections[0]["title"]
        left_content = fallback_string(sections[0]["content"])
        right_title = sections[1]["title"]
        right_content = fallback_string(sections[1]["content"])
        return ComparisonSlide(prs, title, left_title, left_content, right_title, right_content)
    
    def save(self, filename):
        self.prs.save(filename)