from pptx import Presentation
from pptx.util import Inches
from .slide_enum import SlideLayout
from urllib.request import urlretrieve
from utils.common import search_unsplash_images
from utils.common import create_image_stream
class PictureWithCaptionSlide:
    def __init__(self, prs, img_path, title_text, caption):
        layout = prs.slide_layouts[SlideLayout.PICTURE_WITH_CAPTION]  # Picture with Caption layout
        slide = prs.slides.add_slide(layout)
        picture_placeholder = slide.placeholders[1]
        left = picture_placeholder.left
        top = picture_placeholder.top
        # print(left, top)
        height = Inches(4)
        local_img_path = urlretrieve(img_path)[0]
        slide.shapes.add_picture(create_image_stream(img_path), left, top, height=height)
        slide.shapes.title.text = title_text
        slide.placeholders[2].text = caption
    
    def save(self, filename):
        self.prs.save(filename)

    @staticmethod
    def from_json(prs, json_slide):
        title = json_slide["title"]
        sections = json_slide["layout"]["sections"]
        img_path = [x for x in sections if x["type"] == "image"][0]["content"]
        caption = [x for x in sections if x["type"] == "caption" or x["type"] == "text"][0]["content"]
        keywords = 'observability engineering'
        # img_path = "https://img.freepik.com/premium-photo/wallpaper-with-dark-dramatic-gradient-colors-ai-generated_88211-6704.jpg"
        img_path = search_unsplash_images(title)
        # print(img_path, caption)
        return PictureWithCaptionSlide(prs, img_path, title, caption)

