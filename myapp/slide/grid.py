from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from .slide_enum import SlideLayout

class GridSlide:
    def __init__(self, prs, title_text, rows, cols, cell_width, cell_height, spacing):
        layout = prs.slide_layouts[SlideLayout.BLANK]  # Blank layout
        slide = prs.slides.add_slide(layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        
        # Add grid shapes
        start_x, start_y = Inches(0.5), Inches(1.5)
        for row in range(rows):
            for col in range(cols):
                left = start_x + col * (cell_width + spacing)
                top = start_y + row * (cell_height + spacing)
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left, top, cell_width, cell_height  # MSO_SHAPE.RECTANGLE is shape type 1
                )
                shape.text = f"R{row+1}C{col+1}"
    
    def save(self, filename):
        self.prs.save(filename)

# Example Usage
# grid_slide = GridSlide(
#     "Grid Slide Example", 
#     rows=3, 
#     cols=4, 
#     cell_width=Inches(1.5), 
#     cell_height=Inches(1), 
#     spacing=Inches(0.2)
# )
# grid_slide.save("grid_slide.pptx")
