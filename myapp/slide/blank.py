from .slide_enum import SlideLayout
from pptx.util import Inches
class BlankSlide:
    def __init__(self, prs, title):
        blank_slide_layout = prs.slide_layouts[SlideLayout.BLANK]  # 6 is for blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add a title to the blank slide
        title = slide.shapes.title
        if not title:
            title = slide.shapes.add_textbox(Inches(2), Inches(1), Inches(6), Inches(1))
            title.text = "Your Slide Title Here"