from pptx import Presentation
from pptx.util import Inches

class TableSlide:
    def __init__(self, prs, title_text):
        title_only_slide_layout = prs.slide_layouts[5]  # Title Only layout
        slide = prs.slides.add_slide(title_only_slide_layout)

        # Set title
        slide.shapes.title.text = title_text

        # Add Table
        rows, cols = 2, 2
        left, top = Inches(2.0), Inches(2.0)
        width, height = Inches(6.0), Inches(0.8)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(4.0)

        # Write column headings
        table.cell(0, 0).text = 'Foo'
        table.cell(0, 1).text = 'Bar'

        # Write body cells
        table.cell(1, 0).text = 'Baz'
        table.cell(1, 1).text = 'Qux'

    def save(self, filename):
        self.prs.save(filename)

# Example Usage
# table_slide = TableSlide("Adding a Table")
# table_slide.save("table_slide.pptx")
