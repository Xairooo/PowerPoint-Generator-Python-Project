from pptx import Presentation
from .slide_enum import SlideLayout
from utils.common import fallback_string
class TitleAndVerticalTextSlide:
    def __init__(self, prs, title_text, vertical_text):
        layout = prs.slide_layouts[SlideLayout.TITLE_AND_VERTICAL_TEXT]  # Title and Vertical Text layout
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title_text
        slide.placeholders[1].text = fallback_string(vertical_text)

    @staticmethod
    def from_json(prs, json_slide):
        title = json_slide["title"]
        sections = json_slide["layout"]["sections"]
        vertical_text = [section for section in sections if section["type"] == 'vertical_text' or section["type"] == 'text']
        return TitleAndVerticalTextSlide(prs, title, vertical_text)
    
    def save(self, filename):
        self.prs.save(filename)