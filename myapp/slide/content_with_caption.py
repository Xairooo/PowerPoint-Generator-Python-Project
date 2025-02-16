from pptx import Presentation
from .slide_enum import SlideLayout
class ContentWithCaptionSlide:
    def __init__(self, prs, title_text, caption, content_text):
        layout = prs.slide_layouts[SlideLayout.CONTENT_WITH_CAPTION]  # Content with Caption layout
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title_text
        slide.placeholders[1].text = caption
        slide.placeholders[2].text = content_text
    
    def save(self, filename):
        self.prs.save(filename)

    @staticmethod
    def from_json(prs, json_slide):
        title = json_slide["title"]
        sections = json_slide["layout"]['sections']
        return ContentWithCaptionSlide(prs, title, sections[0]["content"], sections[1]["content"])

